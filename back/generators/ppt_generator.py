from pptx import Presentation

from pptx.util import Inches

import re

import os

from generators.equation_renderer import render_equation

# ---------------------------------------------------
# GENERATE POWERPOINT
# ---------------------------------------------------

def generate_ppt(title, content, filename):

    prs = Presentation()

    # ---------------------------------------------------
    # TITLE SLIDE
    # ---------------------------------------------------

    slide_layout = prs.slide_layouts[0]

    slide = prs.slides.add_slide(slide_layout)

    slide.shapes.title.text = title

    slide.placeholders[1].text = "AI Engineering Education Platform"

    # ---------------------------------------------------
    # SPLIT INTO SECTIONS
    # ---------------------------------------------------

    sections = re.split(r"\n## ", content)

    for section in sections:

        section = section.strip()

        if not section:
            continue

        lines = section.split("\n")

        slide_title = lines[0].replace("#", "").strip()

        # ---------------------------------------------------
        # CREATE SLIDE
        # ---------------------------------------------------

        slide_layout = prs.slide_layouts[1]

        slide = prs.slides.add_slide(slide_layout)

        slide.shapes.title.text = slide_title

        body = slide.placeholders[1]

        text_frame = body.text_frame

        first_text = True

        top_position = 1.8

        # ---------------------------------------------------
        # PROCESS EACH LINE
        # ---------------------------------------------------

        for line in lines[1:]:

            line = line.strip()

            if not line:
                continue

            # ---------------------------------------------------
            # BLOCK EQUATION
            # ---------------------------------------------------

            if "$$" in line:
                continue

            # ---------------------------------------------------
            # INLINE EQUATION DETECTION
            # ---------------------------------------------------

            latex_matches = re.findall(r"\$(.*?)\$", line)

            if latex_matches:

                clean_line = re.sub(r"\$(.*?)\$", "", line)

                if clean_line.strip():

                    if first_text:

                        text_frame.text = clean_line

                        first_text = False

                    else:

                        p = text_frame.add_paragraph()

                        p.text = clean_line

                # ---------------------------------------------------
                # ADD EQUATION IMAGE
                # ---------------------------------------------------

                for eq in latex_matches:

                    try:

                        equation_path = render_equation(
                            eq,
                            "outputs/equations"
                        )

                        slide.shapes.add_picture(
                            equation_path,
                            Inches(1),
                            Inches(top_position),
                            width=Inches(4)
                        )

                        top_position += 0.8

                    except:
                        pass

                continue

            # ---------------------------------------------------
            # NORMAL TEXT
            # ---------------------------------------------------

            if first_text:

                text_frame.text = line

                first_text = False

            else:

                p = text_frame.add_paragraph()

                p.text = line

                p.level = 0

        # ---------------------------------------------------
        # FONT SIZE
        # ---------------------------------------------------

        for paragraph in text_frame.paragraphs:

            for run in paragraph.runs:

                run.font.size = Inches(0.18)

    # ---------------------------------------------------
    # SAVE POWERPOINT
    # ---------------------------------------------------

    os.makedirs(
        os.path.dirname(filename),
        exist_ok=True
    )

    prs.save(filename)